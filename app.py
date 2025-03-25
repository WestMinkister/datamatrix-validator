# import 문을 먼저 선언
import os
import streamlit as st
import subprocess
import io
import re
import tempfile
import sys
import platform
import numpy as np
from PIL import Image
import time
import shutil
import base64
from io import BytesIO
import json

# 설정 파일 경로 (절대 경로 사용)
CONFIG_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "datamatrix_config.json")
# 디버그: 설정 파일 경로 정의
print(f"설정 파일 경로: {CONFIG_FILE}")

def load_config():
    """설정 파일에서 구성 불러오기"""
    default_config = {
        "b_range_check": False,
        "b_min_value": 80,
        "b_max_value": 250,
        "i_n_check": True,
        "i_to_n_mapping": {str(i): 10 for i in range(10, 60)}
    }
    
    try:
        if os.path.exists(CONFIG_FILE):
            with open(CONFIG_FILE, 'r') as f:
                return json.load(f)
        else:
            # 설정 파일이 없으면 기본 설정 저장
            with open(CONFIG_FILE, 'w') as f:
                json.dump(default_config, f, indent=2)
            return default_config
    except Exception as e:
        st.error(f"설정 파일 로드 중 오류: {str(e)}")
        return default_config

def save_config(config):
    """설정을 파일에 저장"""
    try:
        # 디버그: 저장 경로 출력
        abs_path = os.path.abspath(CONFIG_FILE)
        st.sidebar.info(f"디버그: 설정 저장 시도 - 파일 경로: {abs_path}")
        st.sidebar.info(f"디버그: 저장할 설정 데이터: {config}")
        
        with open(CONFIG_FILE, 'w') as f:
            json.dump(config, f, indent=2)
        
        # 디버그: 저장 성공 메시지
        st.sidebar.success(f"디버그: 설정이 성공적으로 저장되었습니다 ({time.strftime('%H:%M:%S')})")
        return True
    except Exception as e:
        # 자세한 에러 메시지 표시
        st.sidebar.error(f"디버그: 설정 파일 저장 중 오류: {str(e)}")
        return False

def save_current_config():
    """현재 세션에서 설정 값을 파일로 저장"""
    # 디버그: 함수 호출 및 세션 상태 기록
    st.sidebar.info(f"디버그: save_current_config 호출됨 ({time.strftime('%H:%M:%S')})")
    st.sidebar.info(f"디버그: b_range_check 값: {st.session_state.b_range_check}")
    
    config = {
        "b_range_check": st.session_state.b_range_check,
        "b_min_value": st.session_state.b_min_value,
        "b_max_value": st.session_state.b_max_value,
        "i_n_check": st.session_state.i_n_check,
        "i_to_n_mapping": st.session_state.i_to_n_mapping
    }
    
    result = save_config(config)
    # 디버그: 저장 결과 기록
    st.sidebar.info(f"디버그: 설정 저장 결과: {'성공' if result else '실패'}")
    return result

# 페이지 설정을 가장 먼저 호출해야 함
st.set_page_config(
    page_title="DataMatrix 바코드 검증 도구",
    page_icon="🔍",
    layout="wide",
    initial_sidebar_state="expanded"
)

# 초기 설정 스크립트 실행 (페이지 설정 이후)
try:
    if os.path.exists("init_script.sh"):
        subprocess.run(["bash", "init_script.sh"], check=True)
        st.success("시스템 라이브러리 설치 완료")
except Exception as e:
    st.warning(f"시스템 라이브러리 설치 중 오류 발생: {str(e)}")

# CSS 스타일 적용
st.markdown("""
<style>
    .main .block-container {
        padding-top: 2rem;
        padding-bottom: 2rem;
    }
    h1, h2, h3 {
        margin-bottom: 0.5rem;
    }
    .stProgress > div > div > div > div {
        background-color: #4CAF50;
    }
    .success-box {
        padding: 1rem;
        border-radius: 0.5rem;
        background-color: #F0FFF0;
        border: 1px solid #CCFFCC;
    }
    .warning-box {
        padding: 1rem;
        border-radius: 0.5rem;
        background-color: #FFFAF0;
        border: 1px solid #FAEBD7;
    }
    .error-box {
        padding: 1rem;
        border-radius: 0.5rem;
        background-color: #FFF0F0;
        border: 1px solid #FFCCCC;
    }
</style>
""", unsafe_allow_html=True)

# 각 라이브러리 개별 로드 시도
# OpenCV 로드 시도
try:
    import cv2
    HAVE_CV2 = True
except ImportError:
    HAVE_CV2 = False
    st.warning("OpenCV (cv2) 라이브러리를 불러올 수 없습니다. 이미지 처리 기능이 제한됩니다.")

# pylibdmtx 로드 시도
try:
    from pylibdmtx.pylibdmtx import decode
    HAVE_PYLIBDMTX = True
except ImportError:
    HAVE_PYLIBDMTX = False
    import platform
    if platform.system() == "Windows":
        st.warning("pylibdmtx 라이브러리를 불러올 수 없습니다.")
        st.info("Windows에서 pylibdmtx 설치하기: pip install pylibdmtx 후 libdmtx.dll 파일을 Python 실행 폴더에 복사하세요.")
    else:
        st.warning("pylibdmtx 라이브러리를 불러올 수 없습니다. 바코드 검출 기능을 사용할 수 없습니다.")
        st.info("pylibdmtx 설치를 위해서는 libdmtx 시스템 라이브러리가 필요합니다.")
    
    # 폴백 함수 정의
    def decode(image, **kwargs):
        return []

# pdf2image 로드 시도
try:
    import pdf2image
    HAVE_PDF2IMAGE = True
except ImportError:
    HAVE_PDF2IMAGE = False
    st.warning("pdf2image 라이브러리를 불러올 수 없습니다. PDF 이미지 추출 기능이 제한됩니다.")

# pypdfium2 로드 시도
try:
    import pypdfium2 as pdfium
    HAVE_PDFIUM = True
except ImportError:
    HAVE_PDFIUM = False
    st.warning("pypdfium2 라이브러리를 불러올 수 없습니다. PDF 처리 기능이 제한됩니다.")

# Office 관련 라이브러리 로드 시도
try:
    from openpyxl import load_workbook
    HAVE_OPENPYXL = True
except ImportError:
    HAVE_OPENPYXL = False
    st.warning("openpyxl 라이브러리를 불러올 수 없습니다. Excel 파일 처리 기능을 사용할 수 없습니다.")

try:
    from pptx import Presentation
    HAVE_PPTX = True
except ImportError:
    HAVE_PPTX = False
    st.warning("python-pptx 라이브러리를 불러올 수 없습니다. PowerPoint 파일 처리 기능을 사용할 수 없습니다.")

try:
    from PyPDF2 import PdfReader
    HAVE_PYPDF2 = True
except ImportError:
    HAVE_PYPDF2 = False
    st.warning("PyPDF2 라이브러리를 불러올 수 없습니다. PDF 텍스트 추출 기능을 사용할 수 없습니다.")

# 필요한 라이브러리 설치 확인 메시지
st.info("라이브러리 로드가 완료되었습니다. 일부 라이브러리가 로드되지 않은 경우 해당 기능이 제한될 수 있습니다.")

# 필요한 시스템 패키지 확인 (서버에 미리 설치되어 있어야 함)
def check_system_dependencies():
    """시스템에 필요한 라이브러리가 설치되어 있는지 확인"""
    # 운영체제 확인
    import platform
    
    current_os = platform.system()
    
    # Windows에서는 다른 검사 방법 사용
    if current_os == "Windows":
        # Windows용 검사 코드
        try:
            # pylibdmtx가 로드되었는지만 확인
            if not HAVE_PYLIBDMTX:
                st.warning("pylibdmtx 라이브러리를 불러올 수 없습니다. Windows용 설치 방법을 확인하세요.")
                st.info("Windows에 libdmtx를 설치하려면 https://github.com/dmtx/libdmtx/releases 에서 다운로드하세요.")
            
            # LibreOffice 확인 (Windows 방식)
            import os
            libreoffice_paths = [
                "C:\\Program Files\\LibreOffice\\program\\soffice.exe",
                "C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe"
            ]
            libreoffice_found = any(os.path.exists(path) for path in libreoffice_paths)
            
            if not libreoffice_found:
                st.warning("LibreOffice가 설치되어 있지 않습니다. Office 파일 변환이 작동하지 않을 수 있습니다.")
                st.info("LibreOffice를 https://www.libreoffice.org/download/download/ 에서 다운로드하세요.")
                
        except Exception as e:
            st.warning(f"시스템 확인 중 오류 발생: {str(e)}")
            st.info("일부 기능이 제한될 수 있지만, 앱은 계속 작동합니다.")
        
        return
    
    # Linux/macOS 검사 코드 (기존 코드)
    try:
        # libdmtx 확인
        result = subprocess.run(["pkg-config", "--exists", "libdmtx"],
                               stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        if result.returncode != 0:
            st.warning("libdmtx가 설치되어 있지 않습니다. 바코드 검출이 작동하지 않을 수 있습니다.")
            if current_os == "Darwin":  # macOS
                st.info("macOS에서는 'brew install libdmtx'로 설치할 수 있습니다.")
            else:  # Linux
                st.info("Ubuntu에서는 'sudo apt-get install libdmtx0a libdmtx-dev'로 설치할 수 있습니다.")
    
    except Exception as e:
        st.warning(f"시스템 의존성 확인 중 오류 발생: {str(e)}")
        st.info("이 앱이 정상적으로 작동하려면 libdmtx, libreoffice, poppler-utils가 필요합니다.")

# =========================================================
# 유효성 검증 함수
# =========================================================

def validate_44x44_matrix(data, b_range_check=False, b_min_value=0, b_max_value=9999, i_n_check=False, i_to_n_mapping=None):
    """44x44 매트릭스 데이터 검증 함수"""
    result = {"valid": False, "errors": [], "warnings": [], "data": {}, "pattern_match": False}
    
    # 잘못된 구분자(,) 사용 확인
    if re.search(r'C[A-Za-z0-9]{3},', data) or re.search(r'I\d{2},', data) or re.search(r'W(?:LO|SE),', data):
        result["errors"].append("잘못된 구분자(,)를 사용했습니다. 구분자는 '.'(마침표)여야 합니다.")
        # 어떤 식별자에서 잘못된 구분자를 사용했는지 확인
        if re.search(r'C[A-Za-z0-9]{3},', data):
            result["errors"].append("C 식별자 뒤에 잘못된 구분자(,)를 사용했습니다.")
        if re.search(r'I\d{2},', data):
            result["errors"].append("I 식별자 뒤에 잘못된 구분자(,)를 사용했습니다.")
        if re.search(r'W(?:LO|SE),', data):
            result["errors"].append("W 식별자 뒤에 잘못된 구분자(,)를 사용했습니다.")
        if re.search(r'T\d{2},', data):
            result["errors"].append("T 식별자 뒤에 잘못된 구분자(,)를 사용했습니다.")
        if re.search(r'N\d{3},', data):
            result["errors"].append("N 식별자 뒤에 잘못된 구분자(,)를 사용했습니다.")
        if re.search(r'D\d{8},', data):
            result["errors"].append("D 식별자 뒤에 잘못된 구분자(,)를 사용했습니다.")
        if re.search(r'S\d{3},', data):
            result["errors"].append("S 식별자 뒤에 잘못된 구분자(,)를 사용했습니다.")
        if re.search(r'B[0-9]{120},', data):
            result["errors"].append("B 식별자 뒤에 잘못된 구분자(,)를 사용했습니다.")
        return result

    # 바코드 데이터에서 패턴 확인
    pattern = r'C([A-Za-z0-9]{3})\.I(\d{2})\.W(LO|SE)\.T(\d{2})\.N(\d{3})\.D(\d{8})\.S(\d{3})\.B([0-9]{120})\.'
    match = re.search(pattern, data)
    
    if not match:
        result["errors"].append("44x44 매트릭스 형식이 맞지 않습니다.")
        
        # 개별 패턴 확인으로 어떤 부분이 문제인지 디버깅
        C_match = re.search(r'C([A-Za-z0-9]{3})\.', data)
        if not C_match:
            result["errors"].append("C 식별자를 찾을 수 없거나 형식이 올바르지 않습니다")
        
        I_match = re.search(r'I(\d{2})\.', data)
        if not I_match:
            result["errors"].append("I 식별자를 찾을 수 없거나 형식이 올바르지 않습니다")
        
        W_match = re.search(r'W(LO|SE)\.', data)
        if not W_match:
            result["errors"].append("W 식별자를 찾을 수 없거나 형식이 올바르지 않습니다 (LO 또는 SE 값이어야 함)")
        
        T_match = re.search(r'T(\d{2})\.', data)
        if not T_match:
            result["errors"].append("T 식별자를 찾을 수 없거나 형식이 올바르지 않습니다")
        
        N_match = re.search(r'N(\d{3})\.', data)
        if not N_match:
            result["errors"].append("N 식별자를 찾을 수 없거나 형식이 올바르지 않습니다")
        
        D_match = re.search(r'D(\d{8})\.', data)
        if not D_match:
            result["errors"].append("D 식별자를 찾을 수 없거나 형식이 올바르지 않습니다 (YYYYMMDD 형식)")
        
        S_match = re.search(r'S(\d{3})\.', data)
        if not S_match:
            result["errors"].append("S 식별자를 찾을 수 없거나 형식이 올바르지 않습니다")
        
        B_match = re.search(r'B([0-9]{120})\.', data)
        if not B_match:
            result["errors"].append("B 식별자를 찾을 수 없거나 형식이 올바르지 않습니다 (120자리 숫자)")
        
        return result
    
    C_val, I_val, W_val, T_val, N_val, D_val, S_val, B_val = match.groups()
    
    # 데이터 저장
    result["data"] = {
        "C": C_val,
        "I": I_val,
        "W": W_val,
        "T": T_val,
        "N": N_val,
        "D": D_val,
        "S": S_val,
        "B": B_val
    }
    
    result["pattern_match"] = True
    
    # 추가 검증
    # D: 날짜 형식 검증 (YYYYMMDD)
    try:
        year = int(D_val[:4])
        month = int(D_val[4:6])
        day = int(D_val[6:8])
        
        if not (1900 <= year <= 2100):
            result["errors"].append(f"D 식별자: 연도 범위가 올바르지 않습니다 ({year})")
        if not (1 <= month <= 12):
            result["errors"].append(f"D 식별자: 월 범위가 올바르지 않습니다 ({month})")
        if not (1 <= day <= 31):
            result["errors"].append(f"D 식별자: 일 범위가 올바르지 않습니다 ({day})")
    except ValueError:
        result["errors"].append("D 식별자: 날짜 형식이 올바르지 않습니다")
    
    # B: 숫자 세트 검증
    B_sets = []
    non_zero_sets_count = 0
    
    for i in range(0, len(B_val), 4):
        if i+4 <= len(B_val):
            B_set = B_val[i:i+4]
            B_sets.append(B_set)
            if B_set != '0000':
                non_zero_sets_count += 1
    
    # N: B의 세트 수와 일치하는지 확인
    if int(N_val) != non_zero_sets_count:
        result["errors"].append(f"N 식별자: 값 {N_val}이 B 식별자의 비어있지 않은 세트 수 {non_zero_sets_count}와 일치하지 않습니다")
    
    # I 값에 따른 N 최대값 검증
    if i_n_check and i_to_n_mapping and I_val:
        try:
            i_val_int = int(I_val)
            current_n_value = int(N_val)
            
            # I 값에 해당하는 최대 N 값 찾기
            i_val_str = str(i_val_int)
            if i_val_str in i_to_n_mapping:
                max_n = i_to_n_mapping[i_val_str]
                if current_n_value > max_n:
                    result["errors"].append(f"N 식별자: I{I_val}에 대한 N 값이 최대 허용치({max_n})를 초과했습니다 (현재 값: {current_n_value})")
        except (ValueError, TypeError):
            # I 값이나 N 값이 정수로 변환할 수 없는 경우
            pass
    
    # B 값 범위 검사 (활성화된 경우)
    if b_range_check:
        out_of_range_sets = []
        for B_set in B_sets:
            if B_set != '0000':
                b_val = int(B_set)
                if b_val < b_min_value or b_val > b_max_value:
                    out_of_range_sets.append(f"{B_set} ({b_val})")
        
        if out_of_range_sets:
            error_msg = f"B 식별자: 다음 값들이 지정된 범위({b_min_value}~{b_max_value})를 벗어납니다: {', '.join(out_of_range_sets)}"
            result["errors"].append(error_msg)
    
    # B: 숫자 세트가 오름차순인지 및 큰 점프가 있는지 확인
    prev_set = None
    for B_set in B_sets:
        if B_set != '0000':
            if prev_set:
                # 오름차순 확인
                if int(B_set) <= int(prev_set):
                    result["errors"].append(f"B 식별자: 숫자 세트가 오름차순이 아닙니다 ({prev_set} -> {B_set})")
                # 큰 점프 확인 (100 초과)
                elif int(B_set) - int(prev_set) > 100:
                    result["warnings"].append(f"B 식별자: 숫자 세트 간에 큰 점프가 있습니다 ({prev_set} -> {B_set}, 차이: {int(B_set) - int(prev_set)})")
                # 숫자 하나를 건너뛬어도 경고 표시
                elif int(B_set) - int(prev_set) > 1:
                    result["warnings"].append(f"B 식별자: 숫자 세트 간에 순차가 건너뛬어졌습니다 ({prev_set} -> {B_set}, 누락 값: {int(B_set) - int(prev_set) - 1}개)")
            prev_set = B_set
    
    result["valid"] = len(result["errors"]) == 0
    result["has_warnings"] = len(result["warnings"]) > 0
    
    return result

def validate_18x18_matrix(data):
    """18x18 매트릭스 데이터 검증 함수"""
    result = {"valid": False, "errors": [], "data": {}, "pattern_match": False}
    
    # 잘못된 구분자(,) 사용 확인
    if re.search(r'M[A-Za-z0-9]{4},', data) or re.search(r'I\d{2},', data) or re.search(r'C[A-Za-z0-9]{3},', data):
        result["errors"].append("잘못된 구분자(,)를 사용했습니다. 구분자는 '.'(마침표)여야 합니다.")
        # 어떤 식별자에서 잘못된 구분자를 사용했는지 확인
        if re.search(r'M[A-Za-z0-9]{4},', data):
            result["errors"].append("M 식별자 뒤에 잘못된 구분자(,)를 사용했습니다.")
        if re.search(r'I\d{2},', data):
            result["errors"].append("I 식별자 뒤에 잘못된 구분자(,)를 사용했습니다.")
        if re.search(r'C[A-Za-z0-9]{3},', data):
            result["errors"].append("C 식별자 뒤에 잘못된 구분자(,)를 사용했습니다.")
        if re.search(r'P\d{3},', data):
            result["errors"].append("P 식별자 뒤에 잘못된 구분자(,)를 사용했습니다.")
        return result

    # 바코드 데이터에서 패턴 확인
    pattern = r'M([A-Za-z0-9]{4})\.I(\d{2})\.C([A-Za-z0-9]{3})\.P(\d{3})\.'
    match = re.search(pattern, data)
    
    if not match:
        result["errors"].append("18x18 매트릭스 형식이 맞지 않습니다.")
        
        # 개별 패턴 확인으로 어떤 부분이 문제인지 디버깅
        M_match = re.search(r'M([A-Za-z0-9]{4})\.', data)
        if not M_match:
            result["errors"].append("M 식별자를 찾을 수 없거나 형식이 올바르지 않습니다 (4자리 문자+숫자 조합)")
        
        I_match = re.search(r'I(\d{2})\.', data)
        if not I_match:
            result["errors"].append("I 식별자를 찾을 수 없거나 형식이 올바르지 않습니다 (2자리 숫자)")
        
        C_match = re.search(r'C([A-Za-z0-9]{3})\.', data)
        if not C_match:
            result["errors"].append("C 식별자를 찾을 수 없거나 형식이 올바르지 않습니다 (3자리 문자+숫자 조합)")
        
        P_match = re.search(r'P(\d{3})\.', data)
        if not P_match:
            result["errors"].append("P 식별자를 찾을 수 없거나 형식이 올바르지 않습니다 (3자리 숫자)")
        
        return result
    
    M_val, I_val, C_val, P_val = match.groups()
    
    # 데이터 저장
    result["data"] = {
        "M": M_val,
        "I": I_val,
        "C": C_val,
        "P": P_val
    }
    
    result["pattern_match"] = True
    
    # 여기에 필요한 추가 검증 로직 추가
    
    result["valid"] = len(result["errors"]) == 0
    
    return result

def cross_validate_matrices(matrix_44x44, matrix_18x18):
    """두 매트릭스 간의 교차 검증"""
    errors = []
    
    # 둘 중 하나라도 패턴 매치가 실패한 경우
    if not matrix_44x44["pattern_match"] or not matrix_18x18["pattern_match"]:
        return ["교차 검증을 수행할 수 없습니다. 두 매트릭스의 기본 형식이 올바르지 않습니다."]
    
    # 1. [18x18]의 I 값과 [44x44]의 I 값이 동일한지 확인
    if matrix_18x18["data"]["I"] != matrix_44x44["data"]["I"]:
        errors.append(f"교차 검증 실패: [18x18]의 I({matrix_18x18['data']['I']})와 [44x44]의 I({matrix_44x44['data']['I']})가 일치하지 않습니다.")
    
    # 2. [18x18]의 C 값과 [44x44]의 C 값이 일치하는지 확인
    if matrix_18x18["data"]["C"] != matrix_44x44["data"]["C"]:
        errors.append(f"교차 검증 실패: [18x18]의 C({matrix_18x18['data']['C']})와 [44x44]의 C({matrix_44x44['data']['C']})가 일치하지 않습니다.")
    
    return errors if errors else ["교차 검증이 성공적으로 완료되었습니다."]

# =========================================================
# 이미지 처리 및 바코드 검출 함수
# =========================================================

def split_image_for_detection(image):
    """이미지를 여러 영역으로 분할하여 바코드 인식률 향상"""
    width, height = image.size
    sections = []
    
    # 원본 이미지 추가
    sections.append(image)
    
    # 이미지를 상하좌우로 분할 (4분할)
    half_width = width // 2
    half_height = height // 2
    
    # 좌상단
    sections.append(image.crop((0, 0, half_width, half_height)))
    # 우상단
    sections.append(image.crop((half_width, 0, width, half_height)))
    # 좌하단
    sections.append(image.crop((0, half_height, half_width, height)))
    # 우하단
    sections.append(image.crop((half_width, half_height, width, height)))
    
    # 이미지를 가로로 3등분
    third_height = height // 3
    sections.append(image.crop((0, 0, width, third_height)))
    sections.append(image.crop((0, third_height, width, 2*third_height)))
    sections.append(image.crop((0, 2*third_height, width, height)))
    
    # 이미지를 세로로 3등분
    third_width = width // 3
    sections.append(image.crop((0, 0, third_width, height)))
    sections.append(image.crop((third_width, 0, 2*third_width, height)))
    sections.append(image.crop((2*third_width, 0, width, height)))
    
    return sections

# @st.cache_data 데코레이터 제거 (UnhashableParamError 오류 방지)
def enhance_image_for_detection(image):
    """이미지 전처리를 통해 DataMatrix 인식률 향상 (개선 버전)"""
    # OpenCV로 이미지 처리
    img_array = np.array(image)
    
    results = [image]  # 원본 이미지 포함
    
    # 그레이스케일로 변환
    if len(img_array.shape) == 3:  # 컬러 이미지인 경우
        gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
    else:  # 이미 그레이스케일인 경우
        gray = img_array
    
    # 기본 처리: 노이즈 제거
    denoised = cv2.GaussianBlur(gray, (5, 5), 0)
    
    # 이미지 크기 조정 (확대)
    height, width = gray.shape
    scale_factors = [1.5, 2.0]
    for scale in scale_factors:
        resized = cv2.resize(gray, (int(width * scale), int(height * scale)),
                            interpolation=cv2.INTER_CUBIC)
        results.append(Image.fromarray(resized))
    
    # 여러 이진화 방법 적용
    # 1. 적응형 이진화 (Adaptive Thresholding)
    binary_adaptive = cv2.adaptiveThreshold(denoised, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
                                           cv2.THRESH_BINARY, 11, 2)
    results.append(Image.fromarray(binary_adaptive))
    
    # 2. Otsu 이진화
    _, binary_otsu = cv2.threshold(denoised, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
    results.append(Image.fromarray(binary_otsu))
    
    # 3. 반전된 이진화 (바코드가 역상인 경우)
    _, binary_inv = cv2.threshold(denoised, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU)
    results.append(Image.fromarray(binary_inv))
    
    # 대비 향상 (CLAHE)
    clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
    enhanced = clahe.apply(gray)
    results.append(Image.fromarray(enhanced))
    
    # CLAHE 적용 후 이진화
    _, clahe_binary = cv2.threshold(enhanced, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
    results.append(Image.fromarray(clahe_binary))
    
    # 모폴로지 연산
    kernels = [(3, 3), (5, 5)]
    for k_size in kernels:
        kernel = np.ones(k_size, np.uint8)
        
        # 열림 연산 (침식 후 팽창) - 작은 노이즈 제거
        morph_open = cv2.morphologyEx(binary_adaptive, cv2.MORPH_OPEN, kernel)
        results.append(Image.fromarray(morph_open))
        
        # 닫힘 연산 (팽창 후 침식) - 작은 구멍 채우기
        morph_close = cv2.morphologyEx(binary_adaptive, cv2.MORPH_CLOSE, kernel)
        results.append(Image.fromarray(morph_close))
    
    # 엣지 검출
    edges = cv2.Canny(denoised, 50, 150)
    results.append(Image.fromarray(edges))
    
    # 선명화 필터
    sharpen_kernel = np.array([[-1, -1, -1], [-1, 9, -1], [-1, -1, -1]])
    sharpened = cv2.filter2D(gray, -1, sharpen_kernel)
    results.append(Image.fromarray(sharpened))
    
    return results

def detect_datamatrix(image, progress_callback=None):
    """이미지에서 DataMatrix 바코드 검출 (개선 버전)"""

    # 원본 이미지 전처리
    processed_images = enhance_image_for_detection(image)
    
    all_results = []
    
    # 원본 이미지의 다양한 처리 버전에서 바코드 검출 시도
    for i, img in enumerate(processed_images):
        if progress_callback:
            progress_callback(10 + (i * 30) // len(processed_images))
        try:
            results = decode(img, timeout=5000, max_count=10)
            if results:
                all_results.extend(results)
        except Exception as e:
            st.warning(f"디코딩 중 오류 발생: {str(e)}")
    
    # 이미지가 복잡하거나 바코드가 작을 경우를 위해 이미지 분할 접근
    if len(all_results) < 2:  # 아직 두 개의 바코드를 찾지 못했다면
        # 이미지 분할
        sections = split_image_for_detection(image)
        
        # 각 섹션에 전처리 적용 및 바코드 검출
        for i, section in enumerate(sections):
            if progress_callback:
                progress_callback(50 + (i * 40) // len(sections))
                
            # 섹션 전처리
            section_processed = enhance_image_for_detection(section)
            
            # 처리된 각 섹션에서 바코드 검출
            for img in section_processed:
                try:
                    results = decode(img, timeout=5000, max_count=10)
                    if results:
                        all_results.extend(results)
                except Exception as e:
                    continue  # 에러는 무시하고 계속 진행
    
    # 중복 제거 (바코드 값 기준)
    unique_data = set()
    decoded_data = []
    
    for result in all_results:
        try:
            data = result.data.decode('utf-8', errors='replace')
            if data not in unique_data:
                unique_data.add(data)
                decoded_data.append(data)
        except Exception as e:
            st.warning(f"결과 디코딩 중 오류 발생: {str(e)}")
    
    if progress_callback:
        progress_callback(100)
        
    return decoded_data

# =========================================================
# 파일 처리 함수
# =========================================================

# 수정된 PDF 처리 함수
def extract_images_from_pdf(file_content, progress_callback=None):
    """PDF 파일에서 페이지별 이미지 추출 (오류 방지 기능 추가)"""
    images = []
    
    # 오류 발생 시 표시할 메시지
    error_messages = []
    
    # PDFIUM으로 시도
    if HAVE_PDFIUM:
        try:
            # 임시 파일 생성
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_file:
                temp_file.write(file_content)
                temp_path = temp_file.name
            
            if progress_callback:
                progress_callback(20)
                
            # pypdfium2로 PDF 이미지 추출 (고해상도)
            pdf = pdfium.PdfDocument(temp_path)
            
            total_pages = len(pdf)
            for page_index in range(total_pages):
                if progress_callback:
                    progress_callback(20 + (page_index * 60) // total_pages)
                    
                # 페이지 렌더링 (고해상도로 렌더링하여 바코드 인식률 향상)
                page = pdf[page_index]
                bitmap = page.render(
                    scale=3.0,  # 고해상도로 렌더링
                    rotation=0,
                    crop=(0, 0, 0, 0)
                )
                
                # 이미지 변환
                pil_image = bitmap.to_pil()
                images.append(pil_image)
                
            # 임시 파일 삭제
            os.unlink(temp_path)
            
            if images:
                return images
            else:
                error_messages.append("pypdfium2로 이미지 추출 실패")
        except Exception as e:
            error_messages.append(f"pypdfium2로 PDF 처리 실패: {str(e)}")
    else:
        error_messages.append("pypdfium2 라이브러리가 설치되지 않음")
    
    # pdf2image로 시도
    if HAVE_PDF2IMAGE:
        try:
            st.info("pdf2image로 이미지 추출 시도 중...")
            
            if progress_callback:
                progress_callback(50)
                
            # 임시 디렉토리 생성
            temp_dir = tempfile.mkdtemp()
            temp_pdf_path = os.path.join(temp_dir, 'temp.pdf')
            
            # 파일 저장
            with open(temp_pdf_path, 'wb') as f:
                f.write(file_content)
            
            # pdf2image로 PDF에서 이미지 추출
            pdf_images = pdf2image.convert_from_path(temp_pdf_path, dpi=300)
            images.extend(pdf_images)
            
            # 임시 디렉토리 삭제
            shutil.rmtree(temp_dir)
            
            if progress_callback:
                progress_callback(100)
                
            if images:
                return images
            else:
                error_messages.append("pdf2image로 이미지 추출 실패")
        except Exception as e:
            error_messages.append(f"PDF 파일 처리 실패: {str(e)}")
    else:
        error_messages.append("pdf2image 라이브러리가 설치되지 않음")
    
    # 모든 방법 실패 시
    for msg in error_messages:
        st.error(msg)
    
    st.error("PDF에서 이미지를 추출할 수 없습니다. 필요한 라이브러리가 설치되어 있는지 확인하세요.")
    st.info("PDF 처리를 위해 다음 패키지가 필요합니다: pypdfium2, pdf2image, poppler-utils")
    
    return images

def convert_office_to_pdf(file_content, file_extension, progress_callback=None):
    """Office 파일(PPTX, XLSX)을 PDF로 변환 (LibreOffice 사용)"""
    try:
        if progress_callback:
            progress_callback(10)
            
        # 임시 디렉토리 생성
        temp_dir = tempfile.mkdtemp()
        input_path = os.path.join(temp_dir, f'input.{file_extension}')
        output_path = os.path.join(temp_dir, 'output.pdf')
        
        # 입력 파일 저장
        with open(input_path, 'wb') as f:
            f.write(file_content)
            
        if progress_callback:
            progress_callback(30)
            
        # 운영체제 확인
        import platform
        
        # LibreOffice로 PDF 변환 (OS별 명령어 분기)
        if platform.system() == "Windows":
            # Windows용 명령어
            libreoffice_paths = [
                "C:\\Program Files\\LibreOffice\\program\\soffice.exe",
                "C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe"
            ]
            
            libreoffice_path = None
            for path in libreoffice_paths:
                if os.path.exists(path):
                    libreoffice_path = f'"{path}"'
                    break
                    
            if not libreoffice_path:
                st.warning("LibreOffice를 찾을 수 없습니다.")
                return None
                
            cmd = f'{libreoffice_path} --headless --convert-to pdf --outdir "{temp_dir}" "{input_path}"'
        else:
            # Linux/macOS용 명령어
            cmd = f'libreoffice --headless --convert-to pdf --outdir {temp_dir} {input_path}'
        
        process = subprocess.Popen(cmd, shell=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        _, stderr = process.communicate()
                
        if process.returncode != 0:
            st.warning(f"LibreOffice 변환 실패: {stderr.decode('utf-8', errors='ignore')}")
            return None
            
        if progress_callback:
            progress_callback(70)
            
        # 생성된 PDF 파일 읽기
        try:
            with open(output_path, 'rb') as f:
                pdf_content = f.read()
        except FileNotFoundError:
            st.warning(f"변환된 PDF 파일을 찾을 수 없습니다. LibreOffice가 제대로 설치되어 있는지 확인하세요.")
            return None
        
        # 임시 디렉토리 삭제
        shutil.rmtree(temp_dir)
        
        if progress_callback:
            progress_callback(100)
            
        return pdf_content
    except Exception as e:
        st.error(f"파일 변환 중 오류 발생: {str(e)}")
        return None

def extract_images_from_office_file(file_content, file_extension, progress_callback=None):
    """Office 파일에서 이미지 추출 (PDF 변환 후 처리) - 슬라이드 정보 유지"""
    slide_images = {}  # 슬라이드별 이미지 그룹화
    
    # PDF로 변환
    if progress_callback:
        progress_callback(10, "Office 파일을 PDF로 변환 중...")
        
    pdf_content = convert_office_to_pdf(file_content, file_extension,
                                       lambda p: progress_callback(p * 0.4, "Office 파일을 PDF로 변환 중..."))
    
    if pdf_content:
        # PDF에서 이미지 추출
        if progress_callback:
            progress_callback(50, "PDF에서 이미지 추출 중...")
            
        images = extract_images_from_pdf(pdf_content,
                                       lambda p: progress_callback(50 + p * 0.5, "PDF에서 이미지 추출 중..."))
        
        # 각 이미지를 슬라이드 번호별로 저장
        for i, image in enumerate(images):
            slide_num = i + 1
            if slide_num not in slide_images:
                slide_images[slide_num] = []
            slide_images[slide_num].append(image)
    else:
        st.warning(f"{file_extension.upper()} 파일을 PDF로 변환하지 못했습니다.")
        
        # 직접 이미지 추출 시도 (PPTX만 가능)
        if file_extension.lower() == 'pptx':
            if progress_callback:
                progress_callback(60, "PowerPoint에서 직접 이미지 추출 시도 중...")
                
            try:
                # 임시 파일 생성
                with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
                    temp_file.write(file_content)
                    temp_path = temp_file.name
                
                # 프레젠테이션 열기
                presentation = Presentation(temp_path)
                
                # 슬라이드별로 이미지 추출
                total_slides = len(presentation.slides)
                for slide_idx, slide in enumerate(presentation.slides):
                    if progress_callback:
                        progress_callback(60 + (slide_idx * 40) // total_slides, "슬라이드에서 이미지 추출 중...")
                        
                    slide_num = slide_idx + 1
                    slide_images[slide_num] = []
                    
                    for shape in slide.shapes:
                        if hasattr(shape, 'image'):
                            try:
                                image_bytes = shape.image.blob
                                image = Image.open(io.BytesIO(image_bytes))
                                slide_images[slide_num].append(image)
                            except Exception as e:
                                st.warning(f"이미지 추출 중 오류: {str(e)}")
                
                # 임시 파일 삭제
                os.unlink(temp_path)
            except Exception as e:
                st.error(f"PPTX 직접 처리 중 오류 발생: {str(e)}")
    
    if progress_callback:
        progress_callback(100, "이미지 추출 완료")
        
    return slide_images

# =========================================================
# 결과 출력 함수 - Streamlit UI용으로 변환
# =========================================================

def display_barcode_result(idx, data, result, matrix_type="44x44"):
    """바코드 검증 결과를 Streamlit UI에 표시"""
    # 바코드 데이터 출력 - 앞부분과 뒷부분을 함께 표시
    display_data = data
    
    #if len(data) > 60:
    #    display_data = f"{data[:30]}...{data[-20:]}"
    #else:
    #    display_data = data
    
    st.write(f"**바코드 #{idx+1} 데이터:** {display_data}")
    st.write(f"**{matrix_type} 매트릭스** 형식으로 판단됩니다.")
    
    if result["valid"]:
        if "has_warnings" in result and result["has_warnings"]:
            st.warning("⚠️ 매트릭스 형식은 올바르지만 확인이 필요합니다 (확인 필요 항목 발견)")
        else:
            st.success("✅ 매트릭스 형식이 올바릅니다 (전체 규격 검증 완료)")
        
        if matrix_type == "44x44":
            # 44x44 매트릭스 결과 출력
            col1, col2 = st.columns(2)
            with col1:
                st.write("**추출된 데이터 (1):**")
                st.write(f"* 차종 코드(C): {result['data'].get('C', '없음')}")
                st.write(f"* ITEM 코드(I): {result['data'].get('I', '없음')}")
                st.write(f"* 서열/로트(W): {result['data'].get('W', '없음')}")
                st.write(f"* 부품 차종(T): {result['data'].get('T', '없음')}")
            with col2:
                st.write("**추출된 데이터 (2):**")
                st.write(f"* 적입 수량(N): {result['data'].get('N', '없음')}")
                st.write(f"* 입고 날짜(D): {result['data'].get('D', '없음')}")
                st.write(f"* 서열 번호(S): {result['data'].get('S', '없음')}")
            
            b_data = result['data'].get('B', '')
            if len(b_data) > 30:
                b_display = f"{b_data[:15]}...{b_data[-15:]}"
            else:
                b_display = b_data
            st.write(f"* 커밋 넘버(B): {b_display} ({len(b_data)}자리)")
            
            # B 데이터 세트 표시
            if b_data:
                b_sets = [b_data[i:i+4] for i in range(0, len(b_data), 4) if i+4 <= len(b_data)]
                non_zero_sets = [s for s in b_sets if s != '0000']
                if non_zero_sets:
                    with st.expander("커밋 넘버(B) 세트 상세 보기"):
                        st.write(f"총 {len(b_sets)}개 세트 중 {len(non_zero_sets)}개 유효 세트:")
                        for i, s in enumerate(b_sets):
                            if s != '0000':
                                st.write(f"세트 #{i+1}: **{s}**")
                            else:
                                st.write(f"세트 #{i+1}: {s}")
        else:
            # 18x18 매트릭스 결과 출력
            st.write("**추출된 데이터:**")
            col1, col2 = st.columns(2)
            with col1:
                st.write(f"* 제조사 명(M): {result['data'].get('M', '없음')}")
                st.write(f"* ITEM 코드(I): {result['data'].get('I', '없음')}")
            with col2:
                st.write(f"* 차종 코드(C): {result['data'].get('C', '없음')}")
                st.write(f"* 팔레트번호(P): {result['data'].get('P', '없음')}")
    else:
        if result["pattern_match"]:
            st.warning("❌ 매트릭스의 기본 패턴은 일치하나 추가 검증에 실패했습니다:")
        else:
            st.error("❌ 매트릭스 형식이 맞지 않습니다:")
        for msg in result["errors"]:
            st.write(f"* {msg}")

def display_summary_results(page_results):
    """페이지별 검증 결과 요약 테이블 출력 (Streamlit 버전)"""
    st.markdown("## 📊 페이지별 검증 결과 요약")
    
    # 결과 테이블 데이터 준비
    data = []
    for page_num, result in sorted(page_results.items()):
        # 매트릭스 상태
        matrix_44x44 = "✅ 발견" if result["44x44_found"] else "❌ 없음"
        matrix_18x18 = "✅ 발견" if result["18x18_found"] else "❌ 없음"
        
        # 규격 검증 상태
        if result["44x44_found"] and result["44x44_valid"] and result["18x18_found"] and result["18x18_valid"]:
            if result["has_duplicate_44x44"]:
                validation = "❌ 실패 (44x44 중복)"
            elif result["has_warnings"]:
                validation = "⚠️ 확인 필요"
            else:
                validation = "✅ 통과"
        elif (not result["44x44_found"]) or (not result["18x18_found"]):
            validation = "❌ 실패 (미발견)"
        elif (not result["44x44_valid"]) or (not result["18x18_valid"]):
            validation = "❌ 실패 (규격불일치)"
        else:
            validation = "⚠️ 일부만 통과"
        
        # 교차 검증 상태
        if result["44x44_found"] and result["18x18_found"] and result["cross_valid"]:
            cross_validation = "✅ 통과"
        elif not (result["44x44_found"] and result["18x18_found"]):
            cross_validation = "❓ 검증불가"
        else:
            cross_validation = "❌ 실패"
        
        # 중복 상태
        duplicate_status = "❌ 중복 감지" if result["has_duplicate_44x44"] else "✅ 정상"
            
        data.append([page_num, matrix_44x44, matrix_18x18, validation, cross_validation, duplicate_status])
    
    # 테이블 헤더
    columns = ["페이지/슬라이드", "44x44 검출", "18x18 검출", "규격 검증", "교차 검증", "중복 확인"]
    
    # Streamlit 데이터프레임 표시
    import pandas as pd
    df = pd.DataFrame(data, columns=columns)
    st.dataframe(df, use_container_width=True)
    
    # 최종 결과 출력
    # 경고가 있더라도 결과는 통과로 처리 (경고는 확인 필요 항목일 뿐)
    overall_valid = all(
        result["44x44_found"] and result["44x44_valid"] and
        result["18x18_found"] and result["18x18_valid"] and
        result["cross_valid"] and not result["has_duplicate_44x44"]
        for result in page_results.values()
    )
    
    if overall_valid:
        st.success("✅ 성공: 모든 페이지가 검증을 통과했습니다.")
    else:
        issues_pages = [
            page_num for page_num, result in page_results.items()
            if not (result["44x44_found"] and result["44x44_valid"] and
                   result["18x18_found"] and result["18x18_valid"] and
                   result["cross_valid"]) or result["has_duplicate_44x44"]
        ]
        st.error(f"❌ 실패: {', '.join(map(str, sorted(issues_pages)))} 페이지에서 문제가 발견되었습니다.")

def display_format_help():
    """데이터 매트릭스 형식 정보 출력 (Streamlit 버전)"""
    with st.expander("바코드 형식 안내", expanded=False):
        st.markdown("### 올바른 DataMatrix 바코드 형식")
        
        st.markdown("#### 44x44 매트릭스 형식:")
        st.code("CXXX.IYY.WZZ.TYY.NYYY.DYYYYMMDD.SYYY.BNNNNNNNN...")
        st.markdown("예시: `CAB1.I21.WLO.T10.N010.D20250317.S001.B000100020003000400050006000700080009001000000000000000000000000000000000000000000000000000000000000000000000000000000000.`")
        
        st.markdown("#### 18x18 매트릭스 형식:")
        st.code("MXXXX.IYY.CZZZ.PYYY.")
        st.markdown("예시: `MD213.I30.CSW1.P001.`")
        
        st.markdown("#### 요구사항:")
        req_col1, req_col2 = st.columns(2)
        with req_col1:
            st.markdown("1. 각 페이지에는 두 종류(44x44, 18x18)의 데이터 매트릭스가 있어야 합니다.")
            st.markdown("2. 44x44 매트릭스에서 C, I 식별자 값은 3자리의 (문자+숫자) 조합이어야 합니다.")
            st.markdown("3. 44x44 매트릭스에서 W 식별자 값은 'LO' 또는 'SE'이어야 합니다.")
            st.markdown("4. 44x44 매트릭스의 N 값과 B 식별자의 세트 수가 일치해야 합니다.")
        with req_col2:
            st.markdown("5. 44x44 매트릭스의 B 식별자 숫자 세트는 오름차순이어야 합니다.")
            st.markdown("6. 18x18 매트릭스의 C 값과 44x44 매트릭스의 C 값이 일치해야 합니다.")
            st.markdown("7. 18x18 매트릭스의 I 값과 44x44 매트릭스의 I 값이 일치해야 합니다.")
            st.markdown("8. 모든 식별자는 '.'(마침표)로 구분되어야 합니다.")

# =========================================================
# Streamlit UI 부분
# =========================================================

def main():
    # 설정 파일에서 구성 로드
    config = load_config()
    
    # 디버그: 로드된 설정 출력
    st.sidebar.info(f"디버그: 로드된 설정: {config}")
    
    # 세션 상태 초기화
    if 'admin_mode' not in st.session_state:
        st.session_state.admin_mode = False
    if 'admin_password' not in st.session_state:
        st.session_state.admin_password = "datamatrix_admin"
    if 'b_range_check' not in st.session_state:
        st.session_state.b_range_check = config["b_range_check"]
    if 'b_min_value' not in st.session_state:
        st.session_state.b_min_value = config["b_min_value"]
    if 'b_max_value' not in st.session_state:
        st.session_state.b_max_value = config["b_max_value"]
    if 'i_n_check' not in st.session_state:
        st.session_state.i_n_check = config["i_n_check"]
    if 'i_to_n_mapping' not in st.session_state:
        st.session_state.i_to_n_mapping = config["i_to_n_mapping"]

    # 메인 페이지
    st.title("DataMatrix 바코드 검증 도구 🔍")
    st.markdown("PDF, PowerPoint, Excel 파일에서 DataMatrix 바코드를 검색하고 검증합니다.")
    
    # 보안 배너 추가 (컨테이너 사용)
    security_container = st.container()
    with security_container:
        # 컬럼 비율 조정 (이미지 영역 더 넓게)
        cols = st.columns([1, 3])
        
        with cols[0]:
            # 이미지 크기 확대
            try:
                restricted_img = Image.open("restricted.png")
                st.image(restricted_img, width=200, use_column_width=True)
            except:
                # 이미지가 없는 경우 - 크기 키운 텍스트로 대체
                st.markdown("""
                <div style="font-size: 28px; color: #dc3545; text-align: center; margin-top: 20px;">
                    ⚠️<br>
                    <span style="font-weight: bold;">사내한<br>Restricted</span>
                </div>
                """, unsafe_allow_html=True)
        
        with cols[1]:
            # 텍스트 크기와 마진 조정하여 이미지와 높이 맞추기
            st.markdown("""
            <div style="margin-top: 25px;">
                <h3 style="margin-bottom: 15px;">이 홈페이지는 현대자동차-기아의 정보자산으로 관련 법령에 의해 보호받습니다.</h3>
                <p style="font-size: 18px;">무단 배포 및 복제를 금지하니 취급에 주의하시기 바랍니다.</p>
            </div>
            """, unsafe_allow_html=True)
    
    st.markdown("---")  # 구분선
    st.markdown("PDF, PowerPoint, Excel 파일에서 DataMatrix 바코드를 검색하고 검증합니다.")

    # 사이드바 설정
    with st.sidebar:
        if platform.system() == "Windows":
            st.markdown("---")
        st.markdown("### 관리자 설정")
        
        # 관리자 모드 로그인 UI
        if not st.session_state.admin_mode:
            admin_password = st.text_input("관리자 비밀번호", type="password")
            if st.button("관리자 모드 접속"):
                if admin_password == st.session_state.admin_password:
                    st.session_state.admin_mode = True
                    st.experimental_rerun()
                else:
                    st.error("비밀번호가 올바르지 않습니다.")
        else:
            if st.button("관리자 모드 종료"):
                # 현재 설정 저장
                save_result = save_current_config()
                if save_result:
                    st.success("설정이 저장되었습니다!")
                
                st.session_state.admin_mode = False
                st.experimental_rerun()
            
            st.success("관리자 모드로 로그인했습니다.")
            st.markdown("### 매트릭스 검증 설정 (관리자 전용)")
        
            # B 식별자 범위 설정 UI
            st.markdown("#### B 식별자 값 범위 설정")
            st.markdown("44x44 매트릭스의 B 식별자 뒤에 오는 값(4자리 세트)의 허용 범위를 설정합니다.")
            
            # 체크박스 상태 관리 개선
            if "b_range_check_key" not in st.session_state:
                st.session_state.b_range_check_key = st.session_state.b_range_check
                
            if st.checkbox(
                "B 식별자 범위 검사 활성화",
                value=st.session_state.b_range_check,
                key="b_range_check_key"
            ):
                # 체크박스가 체크되면 True로 설정
                st.session_state.b_range_check = True
                st.sidebar.info(f"디버그: B 식별자 범위 체크박스 설정됨: {st.session_state.b_range_check}")
            else:
                # 체크박스가 해제되면 False로 설정
                st.session_state.b_range_check = False
                st.sidebar.info(f"디버그: B 식별자 범위 체크박스 해제됨: {st.session_state.b_range_check}")
            
            col1, col2 = st.columns(2)
            with col1:
                # 내 중요한 설정 변경만 하면 on_change를 사용하는 것이 더 좋습니다.
                st.session_state.b_min_value = st.number_input(
                    "최소값",
                    min_value=0,
                    max_value=9999,
                    value=st.session_state.b_min_value,
                    key="b_min_value_key"
                )
            with col2:
                st.session_state.b_max_value = st.number_input(
                    "최대값",
                    min_value=0,
                    max_value=9999,
                    value=st.session_state.b_max_value,
                    key="b_max_value_key"
                )
            
            if st.session_state.b_range_check:
                st.info(f"B 식별자 값이 {st.session_state.b_min_value}~{st.session_state.b_max_value} 범위 내에 있는지 검사합니다.")
            
            # I-N 관계 검증 설정 UI
            st.markdown("#### I 식별자에 따른 N 최대값 설정")
            st.markdown("I 식별자 값에 따라 N 식별자가 가질 수 있는 최대값을 설정합니다.")
            
            # 체크박스 상태 관리 개선
            if "i_n_check_key" not in st.session_state:
                st.session_state.i_n_check_key = st.session_state.i_n_check
                
            if st.checkbox(
                "I-N 관계 검사 활성화",
                value=st.session_state.i_n_check,
                key="i_n_check_key"
            ):
                # 체크박스가 체크되면 True로 설정
                st.session_state.i_n_check = True
                st.sidebar.info(f"디버그: I-N 관계 검사 체크박스 설정됨: {st.session_state.i_n_check}")
            else:
                # 체크박스가 해제되면 False로 설정
                st.session_state.i_n_check = False
                st.sidebar.info(f"디버그: I-N 관계 검사 체크박스 해제됨: {st.session_state.i_n_check}")
            
            if st.session_state.i_n_check:
                st.info("각 I 값에 대한 N 최대값을 설정합니다.")
                
                # I10~I59 값을 5개씩 묶어서 표시
                ranges = [(10, 19), (20, 29), (30, 39), (40, 49), (50, 59)]
                
                for start, end in ranges:
                    st.markdown(f"##### I{start} ~ I{end}")
                    cols = st.columns(5)
                    for i, val in enumerate(range(start, end+1)):
                        i_val = str(val)
                        with cols[i % 5]:
                            st.session_state.i_to_n_mapping[i_val] = st.number_input(
                                f"I{i_val}",
                                min_value=1,
                                max_value=999,
                                value=st.session_state.i_to_n_mapping.get(i_val, 10),
                                key=f"i_val_{i_val}_key"
                            )
            
            # 설정 저장 버튼 추가
            save_settings_col1, _ = st.columns([1, 3])
            with save_settings_col1:
                if st.button("설정 저장", key="save_settings_button"):
                    save_result = save_current_config()
                    if save_result:
                        st.success("설정이 성공적으로 저장되었습니다!")
                    else:
                        st.error("설정 저장 중 오류가 발생했습니다.")
                        
            st.markdown("### Windows 환경 설정")
            st.markdown("""
            1. Python 환경에 pylibdmtx 설치: `pip install pylibdmtx`
            2. [libdmtx DLL](https://github.com/dmtx/libdmtx/releases) 다운로드
            3. libdmtx.dll 파일을 Python 실행 경로 또는 시스템 PATH에 추가
            4. [LibreOffice](https://www.libreoffice.org/download/download/) 설치
            """)

        st.title("DataMatrix 바코드 검증 도구")
        st.markdown("---")
        
        st.markdown("### 지원하는 파일 형식")
        st.markdown("* PDF (.pdf)")
        st.markdown("* PowerPoint (.pptx, .ppt)")
        st.markdown("* Excel (.xlsx, .xls)")
        
        st.markdown("---")
        st.markdown("### 바코드 유형")
        st.markdown("* 44x44 매트릭스")
        st.markdown("* 18x18 매트릭스")
        
        st.markdown("---")
        with st.expander("사용 가이드", expanded=False):
            st.markdown("""
            1. 확인할 파일을 업로드합니다.
            2. 각 페이지 또는 슬라이드에서 DataMatrix 바코드가 검색됩니다.
            3. 바코드 데이터가 올바른 형식인지 검증됩니다.
            4. 두 종류의 바코드가 상호 일치하는지 교차 검증이 수행됩니다.
            5. 요약 결과를 통해 전체 결과를 확인합니다.
            """)
            
        st.markdown("---")
        st.markdown("### 시스템 요구사항")
        st.markdown("이 앱은 다음 시스템 패키지를 필요로 합니다:")
        st.markdown("* libdmtx (바코드 디코딩)")
        st.markdown("* poppler-utils (PDF 처리)")
        st.markdown("* libreoffice (Office 파일 변환)")
        
        st.markdown("---")
        st.markdown("### 개발자 정보")
        st.markdown("버전: 1.0.0")
        st.markdown("업데이트: 2025년 3월 21일")
    
    # 시스템 의존성 확인
    if 'system_checked' not in st.session_state:
        check_system_dependencies()
        st.session_state.system_checked = True
    
    # 파일 업로드 기능
    uploaded_file = st.file_uploader("검증할 파일을 업로드하세요",
                                    type=["pdf", "pptx", "ppt", "xlsx", "xls"],
                                    help="PDF, PowerPoint 또는 Excel 파일을 업로드하세요. 각 페이지에서 바코드가 검색됩니다.")
    
    # 바코드 형식 도움말 표시
    display_format_help()
    
    # 파일이 업로드된 경우 처리
    if uploaded_file is not None:
        # 파일 정보 표시
        file_details = {
            "파일명": uploaded_file.name,
            "파일 유형": uploaded_file.type,
            "파일 크기": f"{uploaded_file.size / 1024:.1f} KB"
        }
        
        st.markdown("### 📄 업로드된 파일 정보")
        for key, value in file_details.items():
            st.write(f"**{key}:** {value}")
        
        # 파일 형식 확인
        file_extension = uploaded_file.name.split('.')[-1].lower()
        
        # 진행 상황 표시를 위한 상태 플레이스홀더
        progress_placeholder = st.empty()
        status_placeholder = st.empty()
        
        with st.spinner("파일 처리 중..."):
            progress_bar = progress_placeholder.progress(0)
            status_placeholder.markdown("파일 내용을 읽는 중...")
            
            # 파일 내용 읽기
            file_content = uploaded_file.getvalue()
            
            # 슬라이드별 이미지 그룹화를 위한 딕셔너리
            slide_images = {}
            
            # 파일 형식에 따라 이미지 추출
            if file_extension == 'pdf':
                status_placeholder.markdown("PDF 파일에서 이미지 추출 중...")
                
                # PDF는 페이지별로 이미지 추출
                images = extract_images_from_pdf(file_content,
                                              lambda p: progress_bar.progress(p))
                for i, image in enumerate(images):
                    slide_num = i + 1
                    slide_images[slide_num] = [image]  # 각 페이지를 개별 리스트로 포장
                
                if images:
                    status_placeholder.markdown(f"PDF에서 {len(images)}개 페이지 추출 완료")
                else:
                    status_placeholder.error("PDF에서 이미지를 추출할 수 없습니다.")
                    
            elif file_extension in ['xlsx', 'xls', 'pptx', 'ppt']:
                filetype_name = {'xlsx': 'Excel', 'xls': 'Excel', 'pptx': 'PowerPoint', 'ppt': 'PowerPoint'}
                status_placeholder.markdown(f"{filetype_name[file_extension]} 파일 처리 중...")
                
                # Office 파일에서 슬라이드별 이미지 추출
                def update_progress(p, status="파일 처리 중..."):
                    progress_bar.progress(p)
                    status_placeholder.markdown(status)
                
                slide_images = extract_images_from_office_file(file_content, file_extension, update_progress)
                total_images = sum(len(images) for images in slide_images.values())
                
                if total_images > 0:
                    status_placeholder.markdown(f"{filetype_name[file_extension]}에서 {len(slide_images)}개 슬라이드, {total_images}개 이미지 추출 완료")
                else:
                    status_placeholder.error(f"{filetype_name[file_extension]} 파일에서 이미지를 추출할 수 없습니다.")
            else:
                status_placeholder.error(f"지원되지 않는 파일 형식: {file_extension}")
                slide_images = {}
            
            # 이미지가 추출되었는지 확인
            if not slide_images:
                status_placeholder.error("이미지를 추출할 수 없습니다. 파일이 올바른지 확인하세요.")
                st.stop()
            
            # 페이지별 결과를 저장할 딕셔너리
            page_results = {}
            
            # 44x44 데이터매트릭스 중복 검사를 위한 추적 딕셔너리
            matrices_44x44_track = {}  # key: 데이터 내용, value: 페이지 번호
            
            # 바코드 처리 섹션 헤더
            st.markdown("### 🔎 바코드 검색 및 검증 결과")
            
            # 각 슬라이드/페이지 분석 결과를 보여줄 탭
            page_tabs = st.tabs([f"페이지 {slide_num}" for slide_num in sorted(slide_images.keys())])
            
            # 각 슬라이드/페이지에서 모든 이미지 처리
            for tab_idx, slide_num in enumerate(sorted(slide_images.keys())):
                images = slide_images[slide_num]
                
                with page_tabs[tab_idx]:
                    st.markdown(f"#### 페이지/슬라이드 {slide_num} 분석")
                    st.write(f"슬라이드에서 추출된 이미지: {len(images)}개")
                    
                    # 이미지 미리보기 (접을 수 있는 영역)
                    with st.expander("이미지 미리보기", expanded=False):
                        for img_idx, image in enumerate(images):
                            st.image(image, caption=f"이미지 #{img_idx+1} ({image.width}x{image.height})", use_column_width=True)
                    
                    # 페이지 결과 초기화
                    page_results[slide_num] = {
                        "44x44_found": False,
                        "18x18_found": False,
                        "44x44_valid": False,
                        "18x18_valid": False,
                        "cross_valid": False,
                        "has_duplicate_44x44": False,  # 44x44 중복 감지 필드
                        "duplicate_page": None,  # 중복이 처음 발견된 페이지 번호
                        "has_warnings": False,  # 경고 상태 표시
                        "warning_messages": []  # 경고 메시지 저장
                    }
                    
                    # 이 슬라이드에서 발견된 모든 바코드 저장
                    all_barcodes = []
                    
                    # 바코드 검출 진행 상태 표시
                    barcode_progress = st.progress(0)
                    barcode_status = st.empty()
                    
                    # 각 이미지에서 바코드 검출 및 통합
                    for img_idx, image in enumerate(images):
                        barcode_status.markdown(f"이미지 #{img_idx+1} 바코드 검색 중...")
                        
                        # 이미지에서 데이터매트릭스 검출
                        start_time = time.time()
                        decoded_data = detect_datamatrix(image, lambda p: barcode_progress.progress(p))
                        end_time = time.time()
                        
                        if decoded_data:
                            barcode_status.markdown(f"이미지 #{img_idx+1}에서 {len(decoded_data)}개 바코드 발견 (검색 시간: {end_time - start_time:.2f}초)")
                            all_barcodes.extend(decoded_data)
                        else:
                            barcode_status.warning(f"이미지 #{img_idx+1}에서 바코드를 찾을 수 없습니다 (검색 시간: {end_time - start_time:.2f}초)")
                    
                    # 진행 상태 플레이스홀더 정리
                    barcode_progress.empty()
                    barcode_status.empty()
                    
                    # 중복 제거
                    all_barcodes = list(set(all_barcodes))
                    
                    if not all_barcodes:
                        st.error(f"페이지/슬라이드 {slide_num}에서 DataMatrix 바코드를 찾을 수 없습니다.")
                        continue
                    
                    st.success(f"페이지/슬라이드 {slide_num}에서 총 {len(all_barcodes)}개의 DataMatrix 바코드를 발견했습니다.")
                    
                    # 바코드 데이터 저장 변수
                    data_44x44 = None
                    data_18x18 = None
                    result_44x44 = {"valid": False, "pattern_match": False}
                    result_18x18 = {"valid": False, "pattern_match": False}
                    
                    # 각 바코드 데이터 처리
                    st.markdown("#### 바코드 데이터 검증")
                    
                    for idx, data in enumerate(all_barcodes):
                        # 44x44 매트릭스 패턴 검사
                        if re.search(r'C[A-Za-z0-9]{3}[.,]I\d{2}[.,]W(?:LO|SE)[.,]', data):
                            # 이미 44x44 데이터가 있는 경우 기존 것이 유효한지 확인하고 결정
                            if data_44x44 is None or not result_44x44["valid"]:
                                result_44x44 = validate_44x44_matrix(
                                    data,
                                    b_range_check=st.session_state.b_range_check,
                                    b_min_value=st.session_state.b_min_value,
                                    b_max_value=st.session_state.b_max_value,
                                    i_n_check=st.session_state.i_n_check,
                                    i_to_n_mapping=st.session_state.i_to_n_mapping
                                )
                                data_44x44 = data
                                
                                # Streamlit UI 에 결과 표시
                                st.markdown("##### 44x44 매트릭스 검증")
                                display_barcode_result(idx, data, result_44x44, "44x44")
                                
                                # 결과 업데이트
                                page_results[slide_num]["44x44_found"] = True
                                page_results[slide_num]["44x44_valid"] = result_44x44["valid"]
                                
                                # 경고 상태 업데이트
                                if "has_warnings" in result_44x44 and result_44x44["has_warnings"]:
                                    page_results[slide_num]["has_warnings"] = True
                                    page_results[slide_num]["warning_messages"].extend(result_44x44["warnings"])
                                    
                                    # 경고 메시지 추가 표시
                                    st.warning("⚠️ 확인 필요:")
                                    for warning in result_44x44["warnings"]:
                                        st.write(f"* {warning}")
                        
                        # 18x18 매트릭스 패턴 검사
                        if re.search(r'M[A-Za-z0-9]{4}\.I\d{2}\.C[A-Za-z0-9]{3}\.', data):
                            # 이미 18x18 데이터가 있는 경우 기존 것이 유효한지 확인하고 결정
                            if data_18x18 is None or not result_18x18["valid"]:
                                result_18x18 = validate_18x18_matrix(data)
                                data_18x18 = data
                                
                                # Streamlit UI 에 결과 표시
                                st.markdown("##### 18x18 매트릭스 검증")
                                display_barcode_result(idx, data, result_18x18, "18x18")
                                
                                # 결과 업데이트
                                page_results[slide_num]["18x18_found"] = True
                                page_results[slide_num]["18x18_valid"] = result_18x18["valid"]
                    
                    # 페이지에 두 종류의 매트릭스가 모두 있는지 확인
                    missing_matrix = []
                    if not page_results[slide_num]["44x44_found"]:
                        missing_matrix.append("44x44 매트릭스")
                    if not page_results[slide_num]["18x18_found"]:
                        missing_matrix.append("18x18 매트릭스")
                    
                    if missing_matrix:
                        st.warning(f"⚠️ 경고: 이 페이지에서 {', '.join(missing_matrix)}를 찾을 수 없습니다!")
                    
                    # 44x44 매트릭스 중복 검사
                    if data_44x44 and page_results[slide_num]["44x44_valid"]:
                        if data_44x44 in matrices_44x44_track:
                            # 중복 발견
                            original_page = matrices_44x44_track[data_44x44]
                            page_results[slide_num]["has_duplicate_44x44"] = True
                            page_results[slide_num]["duplicate_page"] = original_page
                            st.error(f"❌ 중복 오류: 페이지 {original_page}에 있는 44x44 매트릭스와 동일한 데이터입니다.")
                        else:
                            # 처음 발견된 경우 추적 딕셔너리에 추가
                            matrices_44x44_track[data_44x44] = slide_num
                    
                    # 교차 검증 수행
                    st.markdown("##### 교차 검증 결과")
                    if data_44x44 and data_18x18:
                        if result_44x44["pattern_match"] and result_18x18["pattern_match"]:
                            cross_results = cross_validate_matrices(result_44x44, result_18x18)
                            
                            if "교차 검증이 성공적으로 완료되었습니다." in cross_results:
                                st.success(cross_results[0])
                                page_results[slide_num]["cross_valid"] = True
                            else:
                                st.error("교차 검증 실패")
                                for msg in cross_results:
                                    st.warning(f"- {msg}")
                                page_results[slide_num]["cross_valid"] = False
                        else:
                            st.error("교차 검증을 수행할 수 없습니다. 두 매트릭스 모두 기본 형식이 일치해야 합니다.")
                    else:
                        st.error("페이지에 44x44와 18x18 매트릭스가 모두 필요합니다.")
            
            # 진행 상태 표시 제거
            progress_placeholder.empty()
            status_placeholder.empty()
            
            # 모든 페이지 분석 후 결과 요약 출력
            if page_results:
                display_summary_results(page_results)
                
                # 결과 다운로드 기능
                st.markdown("### 📥 분석 결과 다운로드")
                
                # 결과를 텍스트로 변환
                report_text = "# DataMatrix 바코드 검증 결과 보고서\n\n"
                report_text += f"## 파일 정보\n"
                report_text += f"- 파일명: {uploaded_file.name}\n"
                report_text += f"- 파일 크기: {file_details['파일 크기']}\n"
                report_text += f"- 처리 날짜: {time.strftime('%Y-%m-%d %H:%M:%S')}\n\n"
                
                report_text += f"## 검증 결과 요약\n"
                issues_pages = [
                    page_num for page_num, result in page_results.items()
                    if not (result["44x44_found"] and result["44x44_valid"] and
                           result["18x18_found"] and result["18x18_valid"] and
                           result["cross_valid"])
                ]
                
                if issues_pages:
                    report_text += f"- 상태: ❌ 실패\n"
                    report_text += f"- 문제 페이지: {', '.join(map(str, sorted(issues_pages)))}\n\n"
                else:
                    report_text += f"- 상태: ✅ 성공\n"
                    report_text += f"- 모든 페이지가 검증을 통과했습니다.\n\n"
                
                report_text += f"## 페이지별 상세 결과\n"
                for page_num, result in sorted(page_results.items()):
                    report_text += f"### 페이지/슬라이드 {page_num}\n"
                    report_text += f"- 44x44 매트릭스: {'발견' if result['44x44_found'] else '없음'}\n"
                    if result['44x44_found']:
                        report_text += f"  - 유효성: {'통과' if result['44x44_valid'] else '실패'}\n"
                        report_text += f"  - 중복 상태: {'❌ 페이지 ' + str(result['duplicate_page']) + '와 중복' if result['has_duplicate_44x44'] else '✅ 중복 없음'}\n"
                        if result['has_warnings']:
                            report_text += f"  - 경고 상태: ⚠️ 확인 필요\n"
                            for warning in result['warning_messages']:
                                report_text += f"    * {warning}\n"
                    report_text += f"- 18x18 매트릭스: {'발견' if result['18x18_found'] else '없음'}\n"
                    if result['18x18_found']:
                        report_text += f"  - 유효성: {'통과' if result['18x18_valid'] else '실패'}\n"
                    if result['44x44_found'] and result['18x18_found']:
                        report_text += f"- 교차 검증: {'통과' if result['cross_valid'] else '실패'}\n"
                    report_text += "\n"
                
                # 다운로드 버튼
                st.download_button(
                    label="📄 분석 결과 보고서 다운로드",
                    data=report_text,
                    file_name=f"datamatrix_report_{time.strftime('%Y%m%d_%H%M%S')}.txt",
                    mime="text/plain",
                )

if __name__ == "__main__":
    main()
